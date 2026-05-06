"""產出三週課程 PPT(低動機跑班版).

輸出:
  teaching_materials/lesson_week1.pptx — 第一節:認識事件,完成事實句
  teaching_materials/lesson_week2.pptx — 第二節:AI 輔助理解,完成差異句
  teaching_materials/lesson_week3.pptx — 第三節:地景連結,完成省思句 + 同儕檢查
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / 'assets' / 'images'
OUT = ROOT / 'teaching_materials'

# ── 設計 token (對齊專案 design_tokens.css) ───────────────────────────────
CREAM = RGBColor(0xF5, 0xF1, 0xEA)
PAPER = RGBColor(0xEC, 0xE5, 0xD7)
INK   = RGBColor(0x1A, 0x14, 0x10)
DEEP  = RGBColor(0x33, 0x2A, 0x22)
SOFT  = RGBColor(0x6B, 0x5E, 0x52)
MUTED = RGBColor(0x96, 0x8B, 0x7E)
LINE  = RGBColor(0xC9, 0xBE, 0xAE)
ACCENT = RGBColor(0xC9, 0xA0, 0x62)  # 金
ALERT = RGBColor(0xB8, 0x47, 0x41)   # 警示朱
HIGHLIGHT = RGBColor(0xD9, 0xA4, 0x41)  # 標重點黃

EVENT_COLORS = {
    'cepo':     RGBColor(0xC8, 0x70, 0x2F),  # 大港口 — 橙朱
    'dafen':    RGBColor(0x4A, 0x5B, 0x8C),  # 大分 — 靛
    'cikasuan': RGBColor(0x6B, 0x8E, 0x7F),  # 七腳川 — 青苔
    'truku':    RGBColor(0x8C, 0x4A, 0x3A),  # 太魯閣 — 赭
}
EVENT_NAMES = {
    'cepo': '大港口事件 / Cepo\' 戰役',
    'dafen': '大分事件',
    'cikasuan': '七腳川事件',
    'truku': '太魯閣戰役',
}

# 16:9 投影
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_SERIF = '思源宋體'      # 思源宋體 / Noto Serif CJK TC, fallback 系統內建
FONT_SANS  = '思源黑體'
FONT_MONO  = 'Consolas'


# ── 小工具 ───────────────────────────────────────────────────────────────

def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text, *,
              font=FONT_SANS, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
              line_spacing=1.4):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box, tf


def _add_bullets(slide, left, top, width, height, items, *,
                 font=FONT_SANS, size=16, color=DEEP, line_spacing=1.6,
                 bullet_color=ACCENT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(8)
        # bullet 點
        b = p.add_run()
        b.text = '・ '
        b.font.name = font
        b.font.size = Pt(size)
        b.font.bold = True
        b.font.color.rgb = bullet_color
        # 內容
        r = p.add_run()
        r.text = item
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def _add_image(slide, path, left, top, width=None, height=None):
    if not Path(path).exists():
        return None
    if width and height:
        return slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    if width:
        return slide.shapes.add_picture(str(path), left, top, width=width)
    return slide.shapes.add_picture(str(path), left, top, height=height)


def _set_notes(slide, text):
    notes = slide.notes_slide.notes_text_frame
    notes.text = text
    for p in notes.paragraphs:
        for r in p.runs:
            r.font.name = FONT_SANS
            r.font.size = Pt(11)


def _add_eyebrow(slide, text, color=MUTED, top=Inches(0.4)):
    _add_text(slide, Inches(0.7), top, Inches(12), Inches(0.3),
              text, font=FONT_SANS, size=10, color=color, line_spacing=1)


def _add_footer(slide, page, total, week_label):
    _add_text(slide, Inches(0.7), Inches(7.05), Inches(8), Inches(0.3),
              f'縱谷無言 ‧ {week_label}', font=FONT_SANS, size=9, color=MUTED)
    _add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.3),
              f'{page} / {total}', font=FONT_SANS, size=9, color=MUTED, align=PP_ALIGN.RIGHT)


# ── 投影片 builder ───────────────────────────────────────────────────────

class Deck:
    def __init__(self, week_label, week_num):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.blank_layout = self.prs.slide_layouts[6]
        self.week_label = week_label
        self.week_num = week_num
        self.slides_specs = []

    def add_slide(self, builder, *args, **kwargs):
        self.slides_specs.append((builder, args, kwargs))

    def render(self, path):
        total = len(self.slides_specs)
        for i, (builder, args, kwargs) in enumerate(self.slides_specs, 1):
            slide = self.prs.slides.add_slide(self.blank_layout)
            _set_bg(slide, CREAM)
            builder(self, slide, *args, **kwargs)
            # footer (cover 與 closing 略過)
            if 1 < i < total:
                _add_footer(slide, i, total, self.week_label)
        self.prs.save(str(path))
        return total


# ── slide layouts (each 是一個函式) ───────────────────────────────────────

def slide_cover(deck, slide, *, title, subtitle, hero, week_text, kicker):
    # 底圖鋪滿(opacity 透過淺色覆蓋層模擬)
    pic = _add_image(slide, hero, Inches(0), Inches(0),
                     width=SLIDE_W, height=SLIDE_H)
    # 半透明卡其覆蓋
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H,
              RGBColor(0xF5, 0xF1, 0xEA))
    # 對 cover 圖貼透明度技巧麻煩,改用半高度貼圖
    if pic:
        slide.shapes._spTree.remove(pic._element)
    # 上半留白 / 下半貼圖
    _add_image(slide, hero, Inches(0), Inches(3.5),
               width=SLIDE_W, height=Inches(4.0))
    # 標題層
    _add_text(slide, Inches(0.8), Inches(0.8), Inches(12), Inches(0.4),
              kicker, font=FONT_SANS, size=11, color=MUTED, line_spacing=1)
    _add_text(slide, Inches(0.8), Inches(1.2), Inches(12), Inches(0.5),
              week_text, font=FONT_SANS, size=14, color=ACCENT, bold=True, line_spacing=1)
    _add_text(slide, Inches(0.8), Inches(1.8), Inches(12), Inches(1.2),
              title, font=FONT_SERIF, size=44, bold=True, color=INK, line_spacing=1.15)
    _add_text(slide, Inches(0.8), Inches(3.0), Inches(12), Inches(0.5),
              subtitle, font=FONT_SERIF, size=20, color=DEEP, line_spacing=1.3)
    # 金線
    _add_rect(slide, Inches(0.8), Inches(2.85), Inches(1.6), Inches(0.04), ACCENT)


def slide_section(deck, slide, *, title, eyebrow, body, color=ACCENT):
    # 左側色帶
    _add_rect(slide, Inches(0), Inches(0), Inches(0.35), SLIDE_H, color)
    _add_eyebrow(slide, eyebrow)
    _add_text(slide, Inches(0.7), Inches(2.3), Inches(12), Inches(1.4),
              title, font=FONT_SERIF, size=40, bold=True, color=INK, line_spacing=1.2)
    if body:
        _add_text(slide, Inches(0.7), Inches(3.9), Inches(12), Inches(2.5),
                  body, font=FONT_SERIF, size=18, color=DEEP, line_spacing=1.7)


def slide_content(deck, slide, *, title, eyebrow, bullets, image=None, image_h=4.0,
                  event=None):
    color = EVENT_COLORS.get(event, ACCENT)
    _add_eyebrow(slide, eyebrow)
    _add_text(slide, Inches(0.7), Inches(0.8), Inches(12), Inches(0.8),
              title, font=FONT_SERIF, size=28, bold=True, color=INK, line_spacing=1.2)
    # 標題下細線
    _add_rect(slide, Inches(0.7), Inches(1.65), Inches(0.8), Inches(0.04), color)
    if image:
        # 右側貼圖
        _add_image(slide, ASSETS / image, Inches(8.4), Inches(2.1),
                   width=Inches(4.4), height=Inches(image_h))
        bullet_w = Inches(7.4)
    else:
        bullet_w = Inches(12)
    _add_bullets(slide, Inches(0.7), Inches(2.0), bullet_w, Inches(4.8),
                 bullets, size=17, color=DEEP, bullet_color=color)


def slide_timetable(deck, slide, *, title, eyebrow, rows):
    _add_eyebrow(slide, eyebrow)
    _add_text(slide, Inches(0.7), Inches(0.8), Inches(12), Inches(0.8),
              title, font=FONT_SERIF, size=28, bold=True, color=INK)
    _add_rect(slide, Inches(0.7), Inches(1.65), Inches(0.8), Inches(0.04), ACCENT)
    # 表格
    top = Inches(2.0)
    row_h = Inches(0.55)
    # 表頭
    _add_rect(slide, Inches(0.7), top, Inches(12), row_h, PAPER, LINE)
    _add_text(slide, Inches(0.85), top + Inches(0.13), Inches(2.0), row_h,
              '時間', font=FONT_SANS, size=13, bold=True, color=DEEP)
    _add_text(slide, Inches(2.85), top + Inches(0.13), Inches(2.5), row_h,
              '長度', font=FONT_SANS, size=13, bold=True, color=DEEP)
    _add_text(slide, Inches(5.35), top + Inches(0.13), Inches(7.4), row_h,
              '老師做什麼 / 學生做什麼', font=FONT_SANS, size=13, bold=True, color=DEEP)
    # 內容
    for i, (t, dur, what) in enumerate(rows):
        y = top + row_h * (i + 1)
        bg = CREAM if i % 2 == 0 else PAPER
        _add_rect(slide, Inches(0.7), y, Inches(12), row_h, bg, LINE)
        _add_text(slide, Inches(0.85), y + Inches(0.15), Inches(2.0), row_h,
                  t, font=FONT_SANS, size=12, bold=True, color=ACCENT)
        _add_text(slide, Inches(2.85), y + Inches(0.15), Inches(2.5), row_h,
                  dur, font=FONT_SANS, size=12, color=SOFT)
        _add_text(slide, Inches(5.35), y + Inches(0.15), Inches(7.4), row_h,
                  what, font=FONT_SANS, size=12, color=DEEP)


def slide_compare(deck, slide, *, title, eyebrow, left_title, left_items,
                  right_title, right_items, left_color=None, right_color=None):
    left_color = left_color or RGBColor(0x6F, 0x8E, 0x6E)   # 綠 = 好
    right_color = right_color or ALERT                       # 紅 = 不夠好
    _add_eyebrow(slide, eyebrow)
    _add_text(slide, Inches(0.7), Inches(0.8), Inches(12), Inches(0.8),
              title, font=FONT_SERIF, size=28, bold=True, color=INK)
    _add_rect(slide, Inches(0.7), Inches(1.65), Inches(0.8), Inches(0.04), ACCENT)
    # 兩欄
    col_w = Inches(5.9)
    # 左
    _add_rect(slide, Inches(0.7), Inches(2.0), col_w, Inches(0.5), left_color)
    _add_text(slide, Inches(0.9), Inches(2.05), col_w, Inches(0.5),
              left_title, font=FONT_SANS, size=14, bold=True, color=CREAM)
    _add_rect(slide, Inches(0.7), Inches(2.5), col_w, Inches(4.3), CREAM, LINE)
    _add_bullets(slide, Inches(0.95), Inches(2.65), col_w - Inches(0.5), Inches(4.0),
                 left_items, size=14, color=DEEP, bullet_color=left_color)
    # 右
    _add_rect(slide, Inches(6.8), Inches(2.0), col_w, Inches(0.5), right_color)
    _add_text(slide, Inches(7.0), Inches(2.05), col_w, Inches(0.5),
              right_title, font=FONT_SANS, size=14, bold=True, color=CREAM)
    _add_rect(slide, Inches(6.8), Inches(2.5), col_w, Inches(4.3), CREAM, LINE)
    _add_bullets(slide, Inches(7.05), Inches(2.65), col_w - Inches(0.5), Inches(4.0),
                 right_items, size=14, color=DEEP, bullet_color=right_color)


def slide_quote_demo(deck, slide, *, title, eyebrow, prompt_label, prompt_text,
                     answer_label, answer_text, hint=None, event=None):
    color = EVENT_COLORS.get(event, ACCENT)
    _add_eyebrow(slide, eyebrow)
    _add_text(slide, Inches(0.7), Inches(0.8), Inches(12), Inches(0.8),
              title, font=FONT_SERIF, size=26, bold=True, color=INK)
    _add_rect(slide, Inches(0.7), Inches(1.65), Inches(0.8), Inches(0.04), color)
    # 提問框
    _add_rect(slide, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5), color)
    _add_text(slide, Inches(0.9), Inches(2.05), Inches(12), Inches(0.5),
              prompt_label, font=FONT_SANS, size=12, bold=True, color=CREAM)
    _add_rect(slide, Inches(0.7), Inches(2.5), Inches(12), Inches(1.4), CREAM, LINE)
    _add_text(slide, Inches(0.95), Inches(2.65), Inches(11.5), Inches(1.2),
              prompt_text, font=FONT_SERIF, size=15, color=DEEP, line_spacing=1.6)
    # 答案框
    _add_rect(slide, Inches(0.7), Inches(4.05), Inches(12), Inches(0.5), DEEP)
    _add_text(slide, Inches(0.9), Inches(4.1), Inches(12), Inches(0.5),
              answer_label, font=FONT_SANS, size=12, bold=True, color=CREAM)
    _add_rect(slide, Inches(0.7), Inches(4.55), Inches(12), Inches(2.0), PAPER, LINE)
    _add_text(slide, Inches(0.95), Inches(4.7), Inches(11.5), Inches(1.8),
              answer_text, font=FONT_SERIF, size=14, color=DEEP, line_spacing=1.7)
    if hint:
        _add_text(slide, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
                  '⚐ ' + hint, font=FONT_SANS, size=11, color=ALERT, line_spacing=1.3)


def slide_check(deck, slide, *, title, eyebrow, items, footer=None):
    _add_eyebrow(slide, eyebrow)
    _add_text(slide, Inches(0.7), Inches(0.8), Inches(12), Inches(0.8),
              title, font=FONT_SERIF, size=28, bold=True, color=INK)
    _add_rect(slide, Inches(0.7), Inches(1.65), Inches(0.8), Inches(0.04), HIGHLIGHT)
    # 黃底高亮
    _add_rect(slide, Inches(0.7), Inches(2.0), Inches(12), Inches(0.45),
              RGBColor(0xF7, 0xE7, 0xC1))
    _add_text(slide, Inches(0.9), Inches(2.05), Inches(12), Inches(0.45),
              '本節最低完成標準(下課前必須有的)',
              font=FONT_SANS, size=13, bold=True, color=DEEP)
    _add_bullets(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(4.0),
                 items, size=18, color=INK, bullet_color=HIGHLIGHT,
                 line_spacing=1.8)
    if footer:
        _add_text(slide, Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
                  footer, font=FONT_SERIF, size=14, color=ACCENT, line_spacing=1.3)


def slide_stance(deck, slide, *, title, eyebrow, words):
    """立場詞 5 欄表格"""
    _add_eyebrow(slide, eyebrow)
    _add_text(slide, Inches(0.7), Inches(0.7), Inches(12), Inches(0.7),
              title, font=FONT_SERIF, size=24, bold=True, color=INK)
    _add_rect(slide, Inches(0.7), Inches(1.4), Inches(0.8), Inches(0.04), ALERT)
    headers = ['詞彙', '漢／日視角', '族視角', '適用事件', '過度簡化風險']
    cols_w = [1.5, 2.7, 2.7, 1.7, 3.4]
    cols_x = [0.7]
    for w in cols_w[:-1]:
        cols_x.append(cols_x[-1] + w)
    top = Inches(1.7)
    row_h = Inches(0.42)
    # 表頭
    for x, w, h in zip(cols_x, cols_w, headers):
        _add_rect(slide, Inches(x), top, Inches(w), row_h, DEEP)
        _add_text(slide, Inches(x + 0.08), top + Inches(0.08), Inches(w), row_h,
                  h, font=FONT_SANS, size=10, bold=True, color=CREAM)
    # 行
    for i, row in enumerate(words):
        y = top + row_h * (i + 1)
        bg = CREAM if i % 2 == 0 else PAPER
        for j, (x, w) in enumerate(zip(cols_x, cols_w)):
            _add_rect(slide, Inches(x), y, Inches(w), row_h, bg, LINE)
            _add_text(slide, Inches(x + 0.08), y + Inches(0.07), Inches(w), row_h,
                      row[j], font=FONT_SANS, size=9, color=DEEP, line_spacing=1.3)


def slide_close(deck, slide, *, title, lines, next_week=None):
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, DEEP)
    _add_text(slide, Inches(0.8), Inches(0.8), Inches(12), Inches(0.4),
              '縱谷無言 ‧ ' + deck.week_label,
              font=FONT_SANS, size=11, color=MUTED)
    _add_text(slide, Inches(0.8), Inches(1.6), Inches(12), Inches(1.0),
              title, font=FONT_SERIF, size=36, bold=True, color=CREAM)
    _add_rect(slide, Inches(0.8), Inches(2.7), Inches(1.6), Inches(0.04), ACCENT)
    _add_bullets(slide, Inches(0.8), Inches(3.0), Inches(12), Inches(2.5),
                 lines, size=17, color=PAPER, bullet_color=ACCENT, line_spacing=1.7)
    if next_week:
        _add_text(slide, Inches(0.8), Inches(6.2), Inches(12), Inches(0.4),
                  '下週',
                  font=FONT_SANS, size=11, color=ACCENT)
        _add_text(slide, Inches(0.8), Inches(6.5), Inches(12), Inches(0.6),
                  next_week, font=FONT_SERIF, size=18, color=CREAM)


# ── Week 1 ───────────────────────────────────────────────────────────────

def build_week1():
    d = Deck(week_label='第一節 認識事件 ‧ 完成事實句', week_num=1)

    # 1 cover
    d.add_slide(slide_cover,
        kicker='花蓮高商 ‧ 多元文化與文學 ‧ 原住民族單元',
        week_text='第一節 / 共三節',
        title='縱谷無言',
        subtitle='花東縱谷上,我們從不知道發生過的事',
        hero='hero_valley.png')

    # Helper to attach notes
    def _attach_note(text):
        # re-wrap last spec to add note rendering
        builder, args, kwargs = d.slides_specs[-1]
        kwargs['_speaker_notes'] = text
        d.slides_specs[-1] = (builder, args, kwargs)

    _attach_note('第一張的 30 秒 — 不要解釋什麼是「縱谷無言」,讓標題沉下去就好。\n'
                 '提示學生:今天三節課,我們會用 NotebookLM 一起把這四個沒人講過的故事撿回來。\n'
                 '節奏:這張不要超過 60 秒。')

    # 2 三週進度地圖
    d.add_slide(slide_section,
        eyebrow='本單元三週進度',
        title='三週,三句話,一張地圖',
        body='第一週:認識事件 + 寫一句「事實句」 — 先逛、有印象就好\n'
             '第二週:寫下「差異句」 — 官方與族人怎麼說不一樣\n'
             '第三週:寫下「省思句」 — 知道後,你怎麼想\n\n'
             '三句話最後會出現在班級成果網頁與花蓮事件地圖上。',
        color=ACCENT)
    _attach_note('讓學生看到全貌:這三節不是各自獨立,而是「漏斗」。\n'
                 '第一週是「探路」 — 不要求觀點思辨,先有印象、有成就感。\n'
                 '深層的觀點對照、差異分析放在第二、三週 — 別讓學生第一週就有壓力。')

    # 3 為什麼是這四個事件
    d.add_slide(slide_content,
        eyebrow='為什麼選這四個',
        title='花蓮這片土地,不是空白的',
        bullets=[
            '我們腳下的這條縱谷,曾經發生過讓族人「說不出話」的事',
            '不是從前才有,1877 一直延伸到 1933:超過 50 年的衝突軌跡',
            '四個事件分別在縱谷不同地理位置:海口、深山、平原、峽谷',
            '今天的地名、學校、街道、紀念碑 — 都還留著這些事的痕跡',
            '這不是歷史考題,這是我們生活的地方的另一層記憶',
        ],
        image='hero_valley.png',
        image_h=4.5)
    _attach_note('用「我們腳下」拉回花蓮在地感。\n'
                 '若學生問「為什麼不是 228、霧社」— 回答:那些事件也重要,但這四個是「花蓮的事」,從你家走得到。\n'
                 '這張不要超過 90 秒,進入事件介紹。')

    # 4-7 四事件快速介紹
    events_intro = [
        ('cepo', '大港口事件 / Cepo\' 戰役', '1877–1878', '阿美族 ‧ 海岸線', 'event_cepo.png',
         ['清廷夏獻綸與兵勇進入花東縱谷東緣',
          '阿美族 Cepo\' 部落被指控「拒撫」、遭武力鎮壓',
          '事件後設「靜浦」,今日的靜浦國小就蓋在當時的衝突地點上',
          '族人記憶:這不是「事件」,是一場戰爭、一場屠殺']),
        ('dafen', '大分事件', '1914–1933', '布農族 ‧ 深山', 'event_dafen.png',
         ['18 年漫長武裝抵抗,而不是一次性「事件」',
          'Dahu Ali(拉荷阿雷)率族人對抗日警駐在所',
          '日方投入大量「理蕃」資源、推進駐在所、收押槍枝',
          '族人記憶:這是守護獵場與部落的長期戰爭']),
        ('cikasuan', '七腳川事件', '1908', '阿美族 ‧ 平原', 'event_cikasuan.png',
         ['七腳川社原本協助日方擔任「隘勇」',
          '因薪資、槍枝、與日方條件衝突,從合作翻轉為對立',
          '日方發動討伐,部落被「滅社」、土地被沒收、族人被遷移',
          '今日吉安的部分土地仍是當時被沒收後重新分配的']),
        ('truku', '太魯閣戰役', '1914', '太魯閣族 ‧ 峽谷', 'event_truku.png',
         ['日方稱「太魯閣蕃討伐」 — 名稱本身就是立場',
          '佐久間左馬太總督親自率軍入山',
          '太魯閣族多年抵抗、最終全面被收編',
          '戰後部落被強制遷出山區,改變整個族群的居住版圖']),
    ]
    for eid, title, year, place, img, bullets in events_intro:
        d.add_slide(slide_content,
            eyebrow=f'四個事件之一 ‧ {year} ‧ {place}',
            title=title,
            bullets=bullets,
            image=img,
            image_h=4.2,
            event=eid)
        _attach_note(
            f'本張 60-90 秒。重點:讓學生「聽得進」這個事件的核心張力。\n'
            f'不要把所有細節都講完 — 等他們去 NotebookLM 自己問。\n'
            f'唸完 4 個 bullet 即可,最後一個 bullet 是「族人記憶」,語氣慢一點。')

    # 8 你被分配到的事件 + 任務說明
    d.add_slide(slide_content,
        eyebrow='接下來 5 分鐘 ‧ 找到你的位置',
        title='你被分配到的事件 — 看老師的名單',
        bullets=[
            '老師已把全班平均分配到四個事件 — 每事件人數相同',
            '看老師投影 / 名單 / 黑板,對應你的座號,找你的事件',
            '同事件同學是「事件夥伴」 — 可以討論,但「每人寫自己的句子」',
            '找到後做兩件事:① 打開「事件理解卡」 ② 等老師示範 NotebookLM',
            '今天目標很簡單:逛一遍 + 寫一句印象話 — 不要求精雕,先有感覺',
        ])
    _attach_note('分配名單必須在這張之前準備好(投影、黑板、列印單擇一)。\n'
                 '快速指出每事件分到幾人 — 讓學生知道「你不孤單」。\n'
                 '同事件可以討論,但寫自己的句子 — 強調這點。\n'
                 '這張結束後直接進 NotebookLM 介紹(下一張)。')

    # 9 NotebookLM 是什麼
    d.add_slide(slide_content,
        eyebrow='今天主要工具',
        title='NotebookLM:不是 ChatGPT、不會亂講',
        bullets=[
            'Google 的 AI 閱讀助手 — 只回答「老師上傳的 PDF 裡」的內容',
            '老師已經把每個事件的兩本權威資料都上傳好',
            '它會「附頁碼」 — 可以查證、可以引用',
            '不只能問問題 — 還能一鍵產出「簡報、心智圖、音訊摘要」',
            '它跟 ChatGPT 最大不同:可信、可引用、有來源',
        ])
    _attach_note('低動機學生很多用過 ChatGPT,但會被「亂講」激怒。\n'
                 '強調:NotebookLM 答錯機率低很多,因為它只看「給它看的」資料。\n'
                 '預告「不只能問問題」 — 不愛打字的學生,可以直接點 Studio 一鍵產出。\n'
                 '接下來 6 張會手把手教操作。')

    # 10 登入步驟 1 — 用學校 Google 帳號
    d.add_slide(slide_content,
        eyebrow='第一次用 NotebookLM ‧ Step 1 / 6',
        title='打開 Chrome,確認登入「學校 Google 帳號」',
        bullets=[
            '一定要用電腦版瀏覽器 — Chrome 或 Edge,手機 / 平板介面不完整',
            '右上角頭像 → 確認帳號是學校信箱(@flgs.hlc.edu.tw 或學校網域)',
            '不要用個人 Gmail — 學校的 Education Viewer 只認學校網域',
            '若右上角是個人帳號 → 點頭像 → 切換,或開新分頁登入學校帳號',
            '無痕視窗也可以 — 但同樣要用學校帳號登',
        ])
    _attach_note('帳號錯誤是第一次最常見問題 — 90% 的「打不開」都是這個。\n'
                 '老師可以舉手:「右上角是學校帳號對嗎?」逐桌目視確認。\n'
                 '若有學生忘了學校密碼 → 立刻請他到旁邊問資訊組,不要卡進度。')

    # 11 你的角色 — Viewer(放心,弄不壞)
    d.add_slide(slide_content,
        eyebrow='第一次用 NotebookLM ‧ Step 2 / 6',
        title='你的角色是「檢視者(Viewer)」 — 你弄不壞它',
        bullets=[
            '老師已把每個 Notebook 設成「Viewer」權限 — 你是「讀」這本書的人',
            '✓ 可以做:看 PDF、問問題、看引用、產出 Studio(簡報/心智圖/音訊)',
            '✗ 做不到(系統幫你擋):刪 PDF、加 PDF、改老師的設定',
            '所以放心按、放心問 — 按錯按鈕也不會弄壞,最多就是答案不滿意',
            '若連結打不開 → 通常是 Google 帳號錯了(看上一張),不是你的鍋',
        ])
    _attach_note('低動機學生最大恐懼是「按錯被罵」 — 第一句就消除這個恐懼。\n'
                 '「Viewer」= 圖書館借書,你能看、能查、能影印筆記,但不能改書本身。\n'
                 '強調:能做的都安全,不能做的系統會擋。所以「亂試」是被允許的。\n'
                 '這張結束直接進介面導覽。')

    # 12 介面導覽 — 三大區
    d.add_slide(slide_content,
        eyebrow='第一次用 NotebookLM ‧ Step 3 / 6',
        title='介面長這樣 — 認識三大區塊',
        bullets=[
            '【左】Sources — 老師上傳的 PDF 列表(可看、不可刪、不可加)',
            '【中】Chat — 你跟 AI 對話的地方,本週主要在這裡',
            '【右】Studio — 一鍵產出「簡報 / 心智圖 / 音訊摘要 / 大綱」 — 本週可玩玩看',
            '中間最下方有「打字框」— 把問題打進去、按 Enter',
            '左側 Sources 可以點開 PDF 預覽 — 用來查證 AI 答案',
        ])
    _attach_note('如果有實機投影,直接打開介面對照講 — 比文字解釋有效 10 倍。\n'
                 '若無投影 → 至少口頭描述位置(左/中/右),強調「打字框在中間下方」。\n'
                 'Studio 區告訴學生「等等示範會玩到」 — 不愛打字的學生會在這裡得分。')

    # 13 怎麼問第一個問題
    d.add_slide(slide_content,
        eyebrow='第一次用 NotebookLM ‧ Step 4 / 6',
        title='怎麼問? — 打字、按 Enter,就這麼簡單',
        bullets=[
            '中間最下方的打字框,把問題打進去',
            '按 Enter(或點送出按鈕)',
            'AI 會「想」5-10 秒,然後出現答案 — 不要重複按,會卡',
            '不必客氣 — 不用寫「請問」「謝謝」,直接問問題',
            '可以追問 — AI 記得上一輪對話脈絡',
        ])
    _attach_note('低動機學生會「不敢問」或「亂按重送」。\n'
                 '兩個強調:① 等 5-10 秒、② 直接問不用客氣。\n'
                 '答案出來後,讓他們「停下來看」 — 不要急著問下一題,接著教看引用。')

    # 14 看引用 — 黃金動作
    d.add_slide(slide_content,
        eyebrow='第一次用 NotebookLM ‧ Step 5 / 6',
        title='答案後面的小數字 — 是黃金,點下去',
        bullets=[
            '每段話後面會有「¹」「²」「³」 — 那是「來源編號」',
            '點數字 → 右側會跳出 PDF 對應頁,並把那一段「高亮」',
            '這就是「可查證」 — AI 不是亂講,有 PDF 為憑',
            '⚠ 沒有來源編號的句子 → 是 AI 推論,引用時要加註明',
            '寫進事件理解卡時,記得寫頁碼(例:P. 23 第 3 段)',
        ])
    _attach_note('「點數字看引用」是本課程最重要的單一動作 — 必須示範至少 2 次。\n'
                 '示範:刻意點一個來源,讓學生看到 PDF 自動跳頁、那段被高亮。\n'
                 '提醒這就是它跟 ChatGPT 最大的不同 — ChatGPT 沒引用、沒辦法查證。')

    # 15 三條鐵則 — 不要做這些
    d.add_slide(slide_compare,
        eyebrow='第一次用 NotebookLM ‧ Step 6 / 6',
        title='三條鐵則 — 一次記住',
        left_title='✓ 可以做',
        left_items=[
            '重複問、換問法、追問',
            '把答案複製成「自己改寫過」的版本',
            '看引用、查 PDF、寫頁碼',
            '跟事件夥伴討論你的問法',
            '把不懂的詞再問 AI 一次',
        ],
        right_title='✗ 不可以做',
        right_items=[
            '把整段答案直接複製貼到事件理解卡(這是抄)',
            '把 NotebookLM 連結 / Studio 產出轉給校外人',
            '截圖貼到 IG / 個人帳號(課堂教材,僅限校內使用)',
            '進「不是你被分配到」的事件 Notebook(請待在你那本)',
            '答案還沒出來就一直按重送 — 等 5-10 秒,它在想',
        ])
    _attach_note('左欄是鼓勵,右欄是警告 — 但不要用威脅語氣。\n'
                 '「viewer 權限」已經把刪/加 PDF 擋掉了 — 這張右欄是「行為層」的提醒,不是技術層。\n'
                 '右欄第 1 點(複製整段)會在後面「事實句寫法」那張再強化一次。')

    # 16 NotebookLM 操作示範 — 問問題 + 產 Studio
    d.add_slide(slide_quote_demo,
        eyebrow='Live demo ‧ 老師示範',
        title='示範:打開大港口 → 問一題 → 點 Studio 產一張心智圖',
        prompt_label='我問 NotebookLM 的問題(輕量版,適合第一次)',
        prompt_text='請用 3-5 句話告訴我:大港口事件大概發生什麼事?'
                    '我是高中生,沒讀過這段歷史,請淺白一點。',
        answer_label='NotebookLM 通常會這樣回(示意)',
        answer_text='1877 年清廷派兵進入花蓮東緣,與阿美族 Cepo\' 部落爆發武力衝突。\n'
                    '事件後清廷在當地設「靜浦」,今天的靜浦地名與此有關。\n'
                    '對阿美族而言,這是一段重要的部落記憶。\n'
                    '【來源:原民會叢書 P.23-25;國教院補充教材 P.6】',
        hint='答案出來後 → 點右側 Studio → Mind Map,讓學生看到「不打字也能整理一張圖」。',
        event='cepo')
    _attach_note('Live demo 兩段示範(寫在備忘):\n'
                 '【段 1:問問題 — 約 3 分鐘】\n'
                 '1. 開 chrome,登入學校 Google\n'
                 '2. 貼 NotebookLM 連結 → 進入「花商二年級_大港口_2026v1」\n'
                 '3. 把問題貼進 chat,按 Enter\n'
                 '4. 答案出來後,點任一句後面的引用 → 看右側 source 跳出對應頁\n'
                 '5. 強調:這就是「可查證」的意思\n\n'
                 '【段 2:產 Studio — 約 3 分鐘】\n'
                 '6. 右側 Studio → Mind Map(或 Briefing Doc)\n'
                 '7. 等 30-60 秒,出來後讓學生「哇」一下 — 這是本節的「成就感點」\n'
                 '8. 順便提:同一個 Notebook 可以產不同類型,等等可以自己玩\n\n'
                 '兩段加起來不要超過 6-7 分鐘 — 剩下時間讓學生自己玩。\n'
                 '若 Mind Map 太慢 → 改 Briefing Doc(較快)。')

    # 17 Studio 五種輸出 — 自己挑一個玩
    d.add_slide(slide_content,
        eyebrow='不愛打字也有得做 ‧ 任選一個',
        title='Studio:右側按鈕,把 Notebook 變成你要的格式',
        bullets=[
            '心智圖(Mind Map):整本 PDF 的關鍵概念地圖,一眼看完全貌',
            '簡報草稿(Slides / Briefing Doc):幾分鐘有一份「你事件的簡報底」',
            '音訊摘要(Audio Overview):兩個 AI 主持人聊你的事件,像 Podcast',
            '時間軸 / 重點摘要 / FAQ / 學習指南:看你想要哪一種',
            '產出來可以「複製、下載、貼到 Notion / 你的筆記」 — 不用重打',
        ])
    _attach_note('這張是「給不愛讀字的學生一條路」 — Studio 是視覺/聽覺學習者的入口。\n'
                 '提醒:Studio 也只看老師上傳的 PDF,不會亂加內容。\n'
                 '但音訊摘要會「發揮」一些 — 像兩個主持人在聊,語氣較鬆 — 當娛樂可、當引用要小心。\n'
                 '本節作業仍以「事件理解卡 + 一句印象話」為主 — Studio 是加分玩法,不是必做。\n'
                 '若學生產出有趣的 Studio 結果 → 鼓勵他下週帶來分享。')

    # 11 事件理解卡開場
    d.add_slide(slide_content,
        eyebrow='你的工作表',
        title='事件理解卡 — 一頁完成本週進度',
        bullets=[
            '老師已發 HTML / docx 兩種格式,選你習慣的',
            '上半部:1 我的事件、2 事件基礎資料(時間/地點/族群/人物)、3 AI 協助',
            '右側:評分軌道(自評用,期末會看)',
            '下方:可展開「AI process log」(進階學生才需要寫)',
            '本週低空標準:基礎資料 4 欄 + 一句「印象話」 — 其他下週再說',
        ])
    _attach_note('提醒不要在第一節就把整張寫完 — 那是三節份的工作。\n'
                 '快速指出:第 7 區「三句展示稿」內,只寫第一句即可,不必雕。\n'
                 '若學生有玩 Studio → 可順手把 Studio 連結 / 截圖貼到 AI process log,當作加分依據。\n'
                 '這張結束後讓學生開始填基礎資料。')

    # 12 事件基礎資料怎麼填
    d.add_slide(slide_compare,
        eyebrow='填空示範 ‧ 大港口為例',
        title='事件基礎資料 4 欄',
        left_title='✓ 這樣寫好',
        left_items=[
            '時間:1877-1878',
            '地點:花蓮縣豐濱鄉靜浦(秀姑巒溪出海口)',
            '族群:阿美族 / Cepo\' 部落',
            '人物:清廷將領夏獻綸、阿美族部落領導 Komod-Pazik',
        ],
        right_title='✗ 這樣寫不夠',
        right_items=[
            '時間:清朝(太籠統 — 哪一年?)',
            '地點:花蓮(範圍太大,要到鄉鎮)',
            '族群:原住民(NotebookLM 會給你「阿美族」)',
            '人物:不知道(去問 NotebookLM,要寫至少一個名字)',
        ])
    _attach_note('立刻講「不知道」不是答案,「去問 NotebookLM」才是答案。\n'
                 '提醒:這 4 欄一定要填,不能空。空白會被退件。\n'
                 '5 分鐘讓學生完成這 4 欄,巡堂幫忙。')

    # 13 事實句句型 — 軟版:「我今天才知道」
    d.add_slide(slide_compare,
        eyebrow='本週只要寫一句 ‧ 不必精雕',
        title='事實句 — 一句你今天記得的事就好',
        left_title='✓ 這樣就達標',
        left_items=[
            '我今天才知道,1877 年靜浦那邊發生過清廷與阿美族的衝突',
            '原來太魯閣族抵抗了好幾年,1914 年才被日軍全面收編',
            '七腳川社一開始幫日方,後來反目被討伐 — 蠻反差的',
            '布農族在大分一帶,從 1914 年抵抗到 1933 年,真的好久',
        ],
        right_title='✗ 這幾種會被退',
        right_items=[
            '不知道(寫一句也好過空白,真的卡再去問 NotebookLM)',
            '一個原住民事件(沒有時間、地點、誰)',
            '【NotebookLM 整段答案】(這是抄,不是你的話)',
            '日本人很壞 / 清朝很爛(立場太重,本週還不寫立場)',
        ])
    _attach_note('關鍵口訣:「我今天才知道 + 時間/誰/在哪/發生什麼」 — 任挑一個切入。\n'
                 '一句話、20-30 字、像自己的話、不能複製 NotebookLM 整段。\n'
                 '本週「不寫立場」 — 立場詞是第二週的工作,這週只是「先有印象」。\n'
                 '5-7 分鐘讓學生寫完,巡堂審 + 給輕鼓勵 — 不要嚴格挑剔。')

    # 14 巡堂時老師問什麼 — 鼓勵版
    d.add_slide(slide_content,
        eyebrow='給老師的口袋提問',
        title='巡堂時 — 用這 4 個問題,鼓勵學生講話',
        bullets=[
            '「你今天記得最清楚的是什麼?哪一句最有印象?」',
            '「你問了 NotebookLM 什麼?它怎麼回?有沒有讓你意外的?」',
            '「你有試 Studio 嗎?產了什麼?願不願意給我看一下?」',
            '「給你 30 秒,跟旁邊事件夥伴講一下你那個事件」',
        ])
    _attach_note('給老師自己的提示卡。第一週是「鼓勵」不是「逼問」 — 讓學生有被聽見的感覺。\n'
                 '不要批評句子寫得爛 — 寫得爛沒關係,有寫就值得讚。\n'
                 '若有學生產出有趣的 Studio → 拍照、放下週開場分享。')

    # 15 第一節最低完成標準 — 低空標準
    d.add_slide(slide_check,
        eyebrow='下課前 5 分鐘 ‧ 自己檢查',
        title='第一節最低完成標準(門檻很低,不要怕)',
        items=[
            '我知道我被分配到哪個事件 ✓',
            '時間、地點、族群、人物 4 欄都有寫(用 NotebookLM 查就好)',
            '我至少跟 NotebookLM 問過一個問題',
            '我寫了一句「印象話」 — 一句即可,不必精雕、不必很長',
        ],
        footer='達標就可以離開。沒達標 → 下週開頭 5 分鐘補就好,不用緊張。')
    _attach_note('下課前 5 分鐘提醒一次。讓學生自己對照。\n'
                 '這四項是「探路」標準 — 高成就的可以多玩 Studio、多寫,但不要求。\n'
                 '深層觀點分析是第二、三週的事,本週無此要求 — 老師不要在這週就拉高標準。\n'
                 '檔案儲存到指定資料夾,或直接交紙本(老師決定)。')

    # 16 收班 + 預告 — 鼓勵版
    d.add_slide(slide_close,
        title='第一節結束 ‧ 你完成了「探路」',
        lines=[
            '今天你做到了:認識事件、用了 NotebookLM、寫了一句印象話',
            '本週是「先逛、有印象」 — 不要求觀點分析,你達標了',
            '下週才會進入「兩邊怎麼說不一樣」 — 那才是真的好玩的部分',
        ],
        next_week='第二節 ‧ 立場詞與差異句 ‧ 你會發現「招撫」與「滅社」是同一件事的兩種說法')
    _attach_note('收班用。第一週強調「你做到了 / 你達標了」 — 不要說「下週還有更多」會嚇到人。\n'
                 '預告下週的「立場詞」是本單元最有趣的部分 — 用「好玩」「反差」等詞 hype 一下。\n'
                 '若本節有學生 Studio 產出特別有趣 → 公告下週開場會請他/她分享 30 秒。')

    return d


# ── Week 2 ───────────────────────────────────────────────────────────────

def build_week2():
    d = Deck(week_label='第二節 AI 輔助理解 ‧ 完成差異句', week_num=2)

    # Helper
    def _attach_note(text):
        builder, args, kwargs = d.slides_specs[-1]
        kwargs['_speaker_notes'] = text
        d.slides_specs[-1] = (builder, args, kwargs)

    # 1 cover
    d.add_slide(slide_cover,
        kicker='花蓮高商 ‧ 多元文化與文學 ‧ 原住民族單元',
        week_text='第二節 / 共三節',
        title='招撫 還是 滅社?',
        subtitle='同樣一件事,兩邊怎麼講不一樣',
        hero='event_dafen.png')
    _attach_note('開場 30 秒。標題自己會講話 — 等學生反應,再進下一張。')

    # 2 上週回顧
    d.add_slide(slide_section,
        eyebrow='上週快速回顧',
        title='你寫了一句「事實句」',
        body='時間 + 誰 + 在哪 + 發生什麼。\n\n'
             '這週要再加一句:同一件事,「不一樣的人會用不一樣的詞」。\n\n'
             '我們把這個叫做「立場詞」 — 詞語會偏一邊。',
        color=EVENT_COLORS['dafen'])
    _attach_note('快速回顧:不要重新講上週,只要點出「事實句已完成」、「這週再加一句」。\n'
                 '不到 60 秒。')

    # 3 為什麼要看不同說法
    d.add_slide(slide_content,
        eyebrow='本週的核心觀念',
        title='歷史不是「一個答案」 — 看誰在講',
        bullets=[
            '同樣一件事,清廷文獻會寫「撫番」、阿美族口述會寫「戰爭」',
            '同樣一件事,日方檔案會寫「討伐」、布農族傳說會寫「抵抗」',
            '哪個是真的?都是真的 — 但「立場」不一樣',
            '我們要做的,是分辨「這個詞站在哪邊」 — 然後決定怎麼寫自己的句子',
        ])
    _attach_note('這是本單元最重要的觀念:歷史的多視角。\n'
                 '不要說「都對」 — 要說「立場不同,所以詞不同」。\n'
                 '為下一張的立場詞表鋪墊。')

    # 4 立場詞表 1 (招撫 / 平亂 / 滅社 / 理蕃)
    d.add_slide(slide_stance,
        eyebrow='立場詞對照 (1/3)',
        title='詞語會偏一邊 — 4 個經典詞',
        words=[
            ['招撫', '清廷視角:勸服歸順', '族人視角:被武力威脅後接受',
             '大港口、七腳川', '把武力包裝成善意'],
            ['平亂', '清廷/日方視角:回復秩序', '族人視角:鎮壓抵抗',
             '大港口、太魯閣', '把抵抗描述為「亂」,合法化武力'],
            ['滅社', '日方視角:消滅部落單位', '族人視角:家園被毀',
             '七腳川', '把屠殺與遷移簡化為行政詞'],
            ['理蕃', '日方視角:管理「未開化」族群', '族人視角:被歧視性政策統治',
             '太魯閣、大分', '「蕃」字本身就是貶意'],
        ])
    _attach_note('每個詞 30 秒。不要逐字念五欄,挑一兩欄重點唸。\n'
                 '舉例方式:「招撫」就是「我來勸你投降」 — 你會覺得是善意嗎?\n'
                 '讓學生筆記哪一個詞最讓他們驚訝。')

    # 5 立場詞表 2
    d.add_slide(slide_stance,
        eyebrow='立場詞對照 (2/3)',
        title='詞語會偏一邊 — 4 個更隱蔽的詞',
        words=[
            ['歸順', '官方視角:自願服從', '族人視角:武力下的選擇',
             '七腳川、太魯閣', '把脅迫包裝成自願'],
            ['收押槍枝', '日方視角:解除武裝、安定治理', '族人視角:剝奪生存與獵場',
             '大分、太魯閣', '只看到「治安」,不看「失去獵場」'],
            ['駐在所', '日方視角:派出所/治理據點', '族人視角:占領、監視、推進線',
             '大分、太魯閣', '「派出所」聽起來中性,實則是占領'],
            ['討伐 / 征討', '日方視角:正當的軍事行動', '族人視角:侵略',
             '太魯閣、大分', '把侵略包裝成正當'],
        ])
    _attach_note('這四個比較隱蔽 — 學生平常聽不出問題。\n'
                 '「駐在所」尤其值得多停 30 秒:聽起來像「派出所」,但功能不一樣。\n'
                 '可以問:「如果有人在你家門口蓋『駐在所』,你會覺得只是治安嗎?」')

    # 6 立場詞表 3
    d.add_slide(slide_stance,
        eyebrow='立場詞對照 (3/3)',
        title='詞語會偏一邊 — 3 個進階詞',
        words=[
            ['隘勇線', '清/日視角:邊界防衛線', '族人視角:被切割部落、限制移動',
             '七腳川、太魯閣', '只看到「防衛」,忽略部落被切'],
            ['太魯閣蕃討伐', '日方視角:歷史事件名稱', '族人視角:命名本身就帶歧視',
             '太魯閣', '名稱本身就是立場 — 連事件名都偏一邊'],
            ['正名', '族人視角:奪回真實名稱', '官方視角:行政程序',
             '所有事件後續', '反向的詞 — 從族人這邊發力的詞'],
        ])
    _attach_note('「太魯閣蕃討伐」這個事件名稱很關鍵 — 維基百科還這樣寫,連事件命名都帶立場。\n'
                 '「正名」是反向的詞 — 從族人那邊出發,可以給學生一點 hope。')

    # 7 立場詞快速練習
    d.add_slide(slide_quote_demo,
        eyebrow='10 分鐘練習 ‧ 找立場詞',
        title='讀這段 NotebookLM 答案 — 找出 3 個立場詞',
        prompt_label='NotebookLM 對「七腳川事件」的回答(節錄)',
        prompt_text='1908 年,七腳川社因不滿日方薪資與條件,發生反抗。日方為平定亂局,'
                    '出兵討伐,並對相關部落進行招撫。最終七腳川社被滅社,土地被收押,'
                    '族人被遷移至他處。',
        answer_label='你應該找到的立場詞 + 為什麼是立場詞',
        answer_text='「亂局」 — 把抵抗叫亂,站在日方一邊\n'
                    '「討伐」 — 把侵略包裝成正當軍事行動\n'
                    '「招撫」 — 把武力威脅包裝成善意\n'
                    '「滅社」 — 簡化屠殺與遷移為行政詞\n'
                    '「收押」 — 把沒收土地包裝為治安',
        hint='5 個都對。請學生找出至少 3 個。',
        event='cikasuan')
    _attach_note('5-7 分鐘讓學生自己找。\n'
                 '提示:不是要他們抗議 NotebookLM,是要他們「看得出」詞站哪邊。\n'
                 '結束後請 1-2 個學生講他們找到哪幾個。')

    # 8 我改掉 AI 的一個地方 — before
    d.add_slide(slide_quote_demo,
        eyebrow='範例示範 ‧ 老師示範改寫',
        title='示範:把上面那段改成「站中間」',
        prompt_label='AI 原句(立場偏日方)',
        prompt_text='1908 年,七腳川社因不滿日方薪資與條件,發生反抗。'
                    '日方為平定亂局,出兵討伐,並對相關部落進行招撫。'
                    '最終七腳川社被滅社,土地被收押,族人被遷移至他處。',
        answer_label='我的改寫(中性 + 雙視角)',
        answer_text='1908 年,七腳川社因薪資與條件爭議,與日方發生衝突。'
                    '日方視為「亂局」並出兵鎮壓,族人視為保護部落的抵抗。'
                    '結果是部落瓦解、土地被沒收、族人被迫遷移,日方稱為「滅社」處置。',
        hint='規則:把帶立場的詞,加上「日方視為...,族人視為...」並列描述。',
        event='cikasuan')
    _attach_note('關鍵教學動作:示範「並列雙視角」的改寫法。\n'
                 '不是「中立」 — 是「兩邊都寫出來」。\n'
                 '這個技巧會用在後面的差異句。\n'
                 '示範後讓學生試著改自己的事件答案。')

    # 9 改寫練習 — 學生時間
    d.add_slide(slide_content,
        eyebrow='15 分鐘 ‧ 你的改寫',
        title='輪到你 — 拿你的事件,改一段 AI 答案',
        bullets=[
            '回到你的事件理解卡 — 第 4 區「我改掉 AI 一個地方」',
            '拿 NotebookLM 給你的某段答案,找出 1-2 個立場詞',
            '改寫成「並列雙視角」 — 日方視為...,族人視為...',
            '不要改太長 — 一兩句就好',
            '這個改寫會被打分(評分軌道 AI 使用 0-2 分)',
        ])
    _attach_note('15 分鐘最少:5 分鐘讓他們選 AI 段落、5 分鐘改、5 分鐘巡堂修。\n'
                 '巡堂時帶著上一張的「並列雙視角」公式去問:「你看,日方怎麼說?族人怎麼說?」')

    # 10 差異句句型
    d.add_slide(slide_content,
        eyebrow='本週要產出的句子',
        title='差異句 — 兩邊怎麼講不一樣',
        bullets=[
            '句型:官方說 ___,但族人記得的是 ___',
            '差異句不是事實的相反 — 是「同一件事的兩種版本」',
            '20-35 字內,像自己的話',
            '會出現在班級成果網頁的卡片上,當「中段亮點句」',
        ])
    _attach_note('差異句的句型只有一個 — 不要讓學生發明新句型。\n'
                 '「官方說 X,但族人記得 Y」就好 — X 與 Y 不同立場、不同重點。\n'
                 '5 分鐘讓學生寫。')

    # 11 差異句好/不夠好對照
    d.add_slide(slide_compare,
        eyebrow='差異句範例對照',
        title='你的差異句 — 像哪邊?',
        left_title='✓ 兩邊都看見',
        left_items=[
            '官方說七腳川社「滅社」,族人記得的是家園被燒毀、被迫遷移',
            '日方稱「太魯閣蕃討伐」,太魯閣族記得的是長達數月的家園防衛',
            '清廷記為「平定 Cepo\' 番」,阿美族傳承的是 Cepo\' 戰爭與屠殺',
            '日方視大分為「治安問題」,布農族視為 18 年的家園守護',
        ],
        right_title='✗ 只看一邊 / 過度評論',
        right_items=[
            '日本很壞欺負原住民(沒有「兩邊」 — 只有評論)',
            '官方都說謊(這不是差異,是抗議)',
            '【NotebookLM 整段】(這是抄,不是你的差異句)',
            '一樣啊都是打仗(沒看到「詞」的差別)',
        ])
    _attach_note('「日本很壞」這種句子會被退件 — 不是因為錯,是因為沒有展現「看見差異」的能力。\n'
                 '提示:好句子會用兩個動詞:官方「說 / 稱 / 記為」、族人「記得 / 視為 / 傳承」。')

    # 12 同學旁邊互看 30 秒
    d.add_slide(slide_content,
        eyebrow='下課前 5 分鐘 ‧ 30 秒互看',
        title='給隔壁同學看 30 秒 — 一個問題',
        bullets=[
            '把你寫的差異句給隔壁同學看',
            '隔壁同學只回答一個問題:「你看得出『兩邊』嗎?」',
            '看得出 → 沒問題',
            '看不太出 → 給對方一個建議怎麼修(不用幫他改,只給建議)',
        ])
    _attach_note('30 秒互看不是同儕互評(那是下週)。是「快速發現自己句子有沒有兩邊」。\n'
                 '不要花太長 — 1-2 分鐘就夠。\n'
                 '巡堂幫忙釐清。')

    # 13 第二節最低完成標準
    d.add_slide(slide_check,
        eyebrow='下課前 ‧ 自己檢查',
        title='第二節最低完成標準',
        items=[
            '事件理解卡第 4 區「改寫 AI」 — 有寫(不能空)',
            '寫了一個差異句 — 包含「官方說 / 族人記得」兩邊',
            '差異句不是 NotebookLM 的整段答案,是「我自己的話」',
            '我可以指出 NotebookLM 答案中至少 1 個立場詞(口頭問可答)',
        ],
        footer='達標就可以離開。下週我們把第三句寫完,還要把三句連回花蓮的地景。')
    _attach_note('最後一項是口頭驗收 — 巡堂時隨便指一個學生、問一次就好。\n'
                 '不是要考倒,是要確認他們腦袋裡有「立場詞」這個概念。')

    # 14 收班 + 預告
    d.add_slide(slide_close,
        title='第二節結束 ‧ 你看見了詞背後的立場',
        lines=[
            '今天你做到的事:不只是讀 AI 答案,還能挑出哪些詞有立場',
            '差異句已寫好 — 你的卡片上現在有「事實 + 差異」兩句',
            '下週最後一節:我們把這兩句連回花蓮今天的地景,寫第三句「省思句」',
        ],
        next_week='第三節 ‧ 地景連結與省思句 ‧ 走過靜浦國小、太魯閣的時候,你會看見什麼?')
    _attach_note('「你做到的事」直接陳述 — 低動機學生需要被肯定。\n'
                 '預告第三節的「地景連結」是課程的高潮 — 把抽象歷史落回具體花蓮。')

    return d


# ── Week 3 ───────────────────────────────────────────────────────────────

def build_week3():
    d = Deck(week_label='第三節 地景連結 ‧ 完成省思句 + 同儕檢查', week_num=3)

    def _attach_note(text):
        builder, args, kwargs = d.slides_specs[-1]
        kwargs['_speaker_notes'] = text
        d.slides_specs[-1] = (builder, args, kwargs)

    # 1 cover
    d.add_slide(slide_cover,
        kicker='花蓮高商 ‧ 多元文化與文學 ‧ 原住民族單元',
        week_text='第三節 / 共三節',
        title='以後我經過⋯',
        subtitle='從歷史走回今天 ‧ 你會看見什麼',
        hero='event_truku.png')
    _attach_note('開場慢一點。標題刻意留白 — 「以後我經過⋯」是省思句的句型起手。\n'
                 '讓學生猜:這句話後面要接什麼?')

    # 2 兩週回顧 + 本週位置
    d.add_slide(slide_section,
        eyebrow='前兩週你寫了',
        title='事實句 + 差異句 = 兩句已完成',
        body='第一週:時間、誰、發生什麼 — 「事實句」\n'
             '第二週:官方說 / 族人記得 — 「差異句」\n\n'
             '今天最後一句:「以後我經過⋯,會想起⋯」 — 省思句。\n'
             '把抽象歷史,變成你下次經過時會記起的事。',
        color=EVENT_COLORS['truku'])
    _attach_note('快速回顧 — 30 秒。\n'
                 '強調本週的不同:不是再寫一句歷史,是「歷史與現在你的連結」。')

    # 3 為什麼這節要連回地景
    d.add_slide(slide_content,
        eyebrow='本週的核心',
        title='歷史不是過去 — 是你現在腳下的土地',
        bullets=[
            '靜浦國小是現在的學校 — 但它蓋在大港口屠殺的地點上',
            '太魯閣國家公園是觀光景點 — 但它的入口曾經是太魯閣族的家',
            '吉安鄉的部分土地 — 是七腳川社被沒收後重新分配的',
            '玉山國家公園的山屋 — 站在當年布農族抵抗的駐在所遺址附近',
            '我們今天的「日常地景」,都疊在某個族群的記憶之上',
        ])
    _attach_note('用具體地名 — 學生熟悉的地方 — 讓「歷史」變得不抽象。\n'
                 '低動機學生需要這個錨:「這跟我有什麼關係」 → 「跟你住的地方有關」。')

    # 4-7 四事件地景現況
    landscape_slides = [
        ('cepo', '大港口 ‧ 今天的靜浦', 'event_cepo.png',
         ['秀姑巒溪出海口至今仍是阿美族重要的祭場',
          '靜浦國小校門口立有紀念碑,寫著「Cepo\' 戰役」(後改名)',
          '校園內就是當年衝突地點,但多數學生不知道',
          '當地仍有 Cepo\' 後代在進行口傳記錄與紀念活動',
          '走過花 11 線時,你經過的就是這個故事']),
        ('dafen', '大分 ‧ 玉山國家公園深處', 'event_dafen.png',
         ['位於今天玉山國家公園東部、海拔 1300m 的高地',
          '當年的駐在所遺址仍在,登山客可達',
          '布農族目前持續推動「拉荷阿雷文化遺址」的保存',
          '這裡是台灣 18 年武裝抗爭中最重要的歷史現場',
          '地點偏遠 — 也因此記憶被外界遺忘最多']),
        ('cikasuan', '七腳川 ‧ 今天的吉安', 'event_cikasuan.png',
         ['位於花蓮市南方的吉安鄉,當年的部落核心已不存',
          '部分土地後被分配給日本農業移民,造成「吉野村」',
          '光復後地名改為「吉安」,移民痕跡仍在街道命名中',
          '七腳川社後人散居台東、花蓮多地,持續推動正名',
          '吉安市區看似平凡 — 但這片地是「重新被分配過」的']),
        ('truku', '太魯閣 ‧ 今天的國家公園', 'event_truku.png',
         ['1914 戰後,太魯閣族被強制遷移至山下平原地區',
          '今天的太魯閣國家公園,是當年的部落生活領域',
          '族人現多居於秀林、卓溪、萬榮等鄉,「下山」已超過 100 年',
          '近年部落推動「傳統領域」的法律認定',
          '當你走過砂卡礑步道,你走在的是別人「以前的家」']),
    ]
    for eid, title, img, bullets in landscape_slides:
        d.add_slide(slide_content,
            eyebrow='地景現況 ‧ 今天看得到什麼',
            title=title,
            bullets=bullets,
            image=img,
            image_h=4.2,
            event=eid)
        _attach_note(f'每張 60-90 秒。\n'
                     f'重點不是地理知識,是「讓學生連結到自己的生活」。\n'
                     f'問一句:「你有去過嗎?有看過這個嗎?」 — 多數會說有。\n'
                     f'然後說:「下次再去,你會想起這個故事。」 — 這就是省思。')

    # 8 「今天留下什麼」勾選
    d.add_slide(slide_content,
        eyebrow='事件理解卡 ‧ 第 6 區',
        title='今天還留下什麼 — 勾選你看到的',
        bullets=[
            '□ 紀念碑 / 紀念建築',
            '□ 地名(部落名、街名、學校名)',
            '□ 學校 / 公共建築',
            '□ 紀念活動 / 祭儀',
            '□ 觀光景點 / 國家公園',
            '□ 幾乎看不出痕跡',
            '可以複選 — 用 NotebookLM 問:「這個事件今天還留下什麼?」',
        ])
    _attach_note('5 分鐘讓學生勾選並補上 NotebookLM 答案。\n'
                 '提醒:選「幾乎看不出痕跡」也是一個答案 — 那本身就是省思。')

    # 9 省思句三種句型
    d.add_slide(slide_content,
        eyebrow='本週要產出的最後一句',
        title='省思句 — 三選一句型',
        bullets=[
            '「我原本以為 ___,但現在發現 ___」 — 認知翻轉型',
            '「這個地名看起來 ___,但背後其實 ___」 — 表象 vs 真相型',
            '「以後我經過 ___,會想起 ___」 — 未來連結型',
            '20-35 字內,寫你「真實」的感受 — 不要寫得像作文',
        ])
    _attach_note('三種句型不要逐字念 — 直接舉一個自己的例子,讓學生模仿。\n'
                 '例:「以後我經過靜浦國小,會想起這裡曾是 Cepo\' 戰役的衝突地。」\n'
                 '5 分鐘讓學生寫。低動機:接受短句、接受不完美 — 但要寫。')

    # 10 省思句好/不夠好
    d.add_slide(slide_compare,
        eyebrow='省思句範例對照',
        title='你的省思句 — 像哪邊?',
        left_title='✓ 真實 + 具體',
        left_items=[
            '我原本以為大港口只是個觀光地名,現在發現它記錄著一場屠殺',
            '太魯閣國家公園看起來是風景,但它曾是太魯閣族的家',
            '以後我走過砂卡礑步道,會想起這條路曾是別人的歸途',
            '我原本以為「招撫」是和平的詞,現在發現它包裝著武力威脅',
        ],
        right_title='✗ 套作文話 / 太空泛',
        right_items=[
            '我學到很多歷史知識(這是作文,不是省思)',
            '原住民很可憐(評論不是省思 — 也不必憐憫)',
            '我們要尊重原住民文化(口號,沒有具體事件)',
            '不知道(寫一個短句也好)',
        ])
    _attach_note('低動機學生最容易寫「我學到很多」 — 直接退件。\n'
                 '退件方式:「這句太作文。再想一下,有沒有一個具體的詞、地點、或想法的轉變?」\n'
                 '5-7 分鐘巡堂幫忙修。')

    # 11 同儕互評 SOP
    d.add_slide(slide_section,
        eyebrow='接下來 10 分鐘 ‧ 同儕互評',
        title='找隔壁 ‧ 看三句 ‧ 給一條建議',
        body='1. 找隔壁同學交換事件理解卡(或螢幕看)\n'
             '2. 只看「事實 / 差異 / 省思」三句展示稿\n'
             '3. 寫下一條具體建議:你建議他改哪一句、怎麼改、為什麼\n'
             '4. 對方收到建議後,選擇「修改一句」或「保留」 — 都要有理由',
        color=ACCENT)
    _attach_note('「找隔壁」的意思 — 不是分組,就是物理上的隔壁。低動機班分組會花太多時間。\n'
                 '「一條建議」是限制,不是建議多多益善 — 太多反而不會動。')

    # 12 同儕建議好/不夠好
    d.add_slide(slide_compare,
        eyebrow='同儕建議的範例',
        title='你給對方的建議 — 像哪邊?',
        left_title='✓ 具體 + 可改',
        left_items=[
            '你的差異句沒提到「族人怎麼說」 — 建議加上「族人記得的是...」',
            '你的省思句寫「很可憐」 — 建議改成具體的一個詞或地名',
            '你的事實句沒寫年份 — 加一個「1908」就更清楚',
            '你三句都有,但「差異」跟「事實」差不多 — 建議差異句寫立場詞',
        ],
        right_title='✗ 太籠統 / 沒幫助',
        right_items=[
            '寫得很好(沒幫助)',
            '改一下(改什麼?改哪句?)',
            '我看不懂(看不懂什麼?)',
            '沒問題(然後呢?)',
        ])
    _attach_note('低動機同儕互評容易寫「OK」、「沒問題」。要逼出具體。\n'
                 '老師示範一下:「你看,『寫得很好』讓人沒辦法改 — 但『差異句沒提族人』讓人馬上知道改哪裡。」')

    # 13 修改一句
    d.add_slide(slide_content,
        eyebrow='接下來 5 分鐘',
        title='收到建議 ‧ 改一句(或拒絕)',
        bullets=[
            '看完同學給的建議',
            '決定:採納 → 改寫一句、不採納 → 寫一行理由為什麼不改',
            '都可以 — 但要有「決定」,不能空白',
            '記在事件理解卡第 9 區「同儕互評」',
        ])
    _attach_note('「拒絕同儕建議」也是合理選擇 — 但要寫理由。\n'
                 '不要強迫採納 — 要訓練「自己決定」。')

    # 14 上傳成果展示
    d.add_slide(slide_content,
        eyebrow='交件方式',
        title='你的三句 ‧ 會出現在這些地方',
        bullets=[
            '班級成果網頁 (class_showcase.html) — 三句展示卡片',
            '花蓮事件地圖 (event_map_interactive.html) — 省思句作為地圖亮點',
            '匿名顯示 — 班級代號 + 座號,不顯示完整姓名',
            '老師會把所有人的卡片解析後產出網頁,你不用做技術',
            '想看自己的句子在網頁長什麼樣 — 課後問老師要連結',
        ])
    _attach_note('讓學生意識到「我的句子會被別人看到」 — 動機提升的關鍵之一。\n'
                 '匿名顯示讓他們安心 — 不會被學弟妹截圖到本名。')

    # 15 整體回望
    d.add_slide(slide_section,
        eyebrow='三節回望',
        title='你做到的事',
        body='三節課,你選了一個事件、查了 NotebookLM、改了 AI 一段話、辨認了立場詞、'
             '寫了三句話、給了同學建議、收到並決定採不採納回饋。\n\n'
             '這是「批判性閱讀 AI 答案」的整套流程 — 你以後讀任何 AI 內容,都會用得上。',
        color=ACCENT)
    _attach_note('明確列舉「你做到的事」 — 不是「我們學到」,是「你做到」。\n'
                 '低動機學生最後一節最需要被看見 — 完成感比知識點重要。')

    # 16 評量規準提醒
    d.add_slide(slide_check,
        eyebrow='評分規準回顧',
        title='你的卡片會這樣被評分',
        items=[
            '基本事實 0-2 分:時間、地點、族群、人物完整性',
            'AI 使用 0-2 分:有沒有改寫、改寫的理由是否清楚',
            '立場比較 0-2 分:差異句有沒有看見「兩邊」',
            '展示句 0-2 分:三句是否清楚、像自己的話、不複製 AI',
        ],
        footer='滿分 8 分,最低通過:基本事實 + 三句展示稿都有。')
    _attach_note('提醒評分軌道 — 但不要強調「你會被評」。\n'
                 '重點是「你已經做了這四件事」 — 評分只是回饋工具。')

    # 17 收班
    d.add_slide(slide_close,
        title='第三節結束 ‧ 縱谷不再無言',
        lines=[
            '三週前,你不一定知道大港口在哪、七腳川為什麼被滅社、Dahu Ali 是誰',
            '今天,你寫下了三句話,這三句話會被同學看到、會被貼在花蓮地圖上',
            '下次經過靜浦、太魯閣、吉安、玉山 — 你不會再「什麼都看不到」',
            '這就是「縱谷不再無言」的意思 — 因為現在,你可以說。',
        ],
        next_week='後續:你的三句會在班級成果網頁與事件地圖上線 — 老師會公告連結')
    _attach_note('整套課程的收尾 — 慢一點、不要趕。\n'
                 '「下次經過...你不會再什麼都看不到」 — 這句話是課程靈魂,慢慢說。\n'
                 '可以選播 30 秒安靜,讓學生看標題 hero 圖。\n'
                 '結束。')

    return d


# ── 主流程:渲染時把 _speaker_notes 取出寫到 notes ────────────────────────

def _render_with_notes(deck, path):
    """覆寫 Deck.render — 把 _speaker_notes 從 kwargs pop 出來,渲染後寫入 notes。"""
    total = len(deck.slides_specs)
    for i, (builder, args, kwargs) in enumerate(deck.slides_specs, 1):
        notes_text = kwargs.pop('_speaker_notes', None)
        slide = deck.prs.slides.add_slide(deck.blank_layout)
        _set_bg(slide, CREAM)
        builder(deck, slide, *args, **kwargs)
        if 1 < i < total:
            _add_footer(slide, i, total, deck.week_label)
        if notes_text:
            _set_notes(slide, notes_text)
    deck.prs.save(str(path))
    return total


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    decks = [
        ('lesson_week1.pptx', build_week1()),
        ('lesson_week2.pptx', build_week2()),
        ('lesson_week3.pptx', build_week3()),
    ]
    for fname, deck in decks:
        path = OUT / fname
        n = _render_with_notes(deck, path)
        print(f'  {fname}  {n} 張  {path}')
    print('完成。')


if __name__ == '__main__':
    main()
